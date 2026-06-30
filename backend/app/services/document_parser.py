"""
Multi-format document parser.
Supports: .pptx, .docx, .pdf, .md
Each parser returns a list of "chunks" — each chunk has text + metadata (page/slide number, title, etc.).
"""

from __future__ import annotations

import io
import re
from typing import List, Optional
from dataclasses import dataclass, field


@dataclass
class ParsedChunk:
    """A single parsed unit (slide, page, or section) from a document."""
    text: str
    page: int = 0
    title: str = ""
    source_type: str = ""  # "slide", "page", "heading"
    meta: dict = field(default_factory=dict)


def parse_pptx(file_bytes: bytes) -> List[ParsedChunk]:
    """Extract text from every slide, including notes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: List[ParsedChunk] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []

        # Slide title
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
            texts.append(f"## {title}")

        # All shape text
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

        # Speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide.notes_text_frame else ""
            if notes_text:
                texts.append(f"[备注] {notes_text}")

        full_text = "\n\n".join(texts)
        if full_text.strip():
            chunks.append(ParsedChunk(
                text=full_text,
                page=idx,
                title=title or f"Slide {idx}",
                source_type="slide",
                meta={"slide_number": idx},
            ))

    return chunks


def parse_docx(file_bytes: bytes) -> List[ParsedChunk]:
    """Extract structured text from a Word document — chunk by paragraphs, group headings."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    chunks: List[ParsedChunk] = []
    current_page = 1
    current_title = ""
    buffer: List[str] = []

    def flush_buffer():
        nonlocal buffer, current_title
        if buffer:
            text = "\n".join(buffer)
            chunks.append(ParsedChunk(
                text=text,
                page=current_page,
                title=current_title,
                source_type="page",
                meta={},
            ))
            buffer = []
            current_title = ""

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            flush_buffer()
            continue

        style = para.style.name if para.style else ""
        is_heading = style and ("Heading" in style or "heading" in style or "Title" in style)

        if is_heading:
            flush_buffer()
            current_title = text
            buffer.append(f"## {text}")
        else:
            buffer.append(text)

    flush_buffer()
    if not chunks:
        # fallback: one big chunk
        full = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if full:
            chunks.append(ParsedChunk(text=full, page=1, title="", source_type="page"))

    return chunks


def parse_pdf(file_bytes: bytes) -> List[ParsedChunk]:
    """Extract text page by page from PDF."""
    import pdfplumber

    chunks: List[ParsedChunk] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                # Try to infer title from first non-empty line
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                title = lines[0] if lines else ""
                chunks.append(ParsedChunk(
                    text=text.strip(),
                    page=idx,
                    title=title,
                    source_type="page",
                    meta={"page_number": idx},
                ))
    return chunks


def parse_markdown(file_bytes: bytes) -> List[ParsedChunk]:
    """Split Markdown by headings into logical sections."""
    text = file_bytes.decode("utf-8", errors="replace")
    chunks: List[ParsedChunk] = []

    # Split on markdown headings (## and #)
    sections = re.split(r"\n(?=#{1,6}\s)", text)
    page = 1
    for section in sections:
        section = section.strip()
        if not section:
            continue
        # Extract heading
        heading_match = re.match(r"^(#{1,6})\s+(.+)", section)
        title = heading_match.group(2).strip() if heading_match else ""
        chunks.append(ParsedChunk(
            text=section,
            page=page,
            title=title,
            source_type="heading",
            meta={},
        ))
        page += 1
    return chunks


# ---------------------------------------------------------------------------
# Unified entry point
# ---------------------------------------------------------------------------
PARSER_MAP = {
    "pptx": parse_pptx,
    "docx": parse_docx,
    "pdf": parse_pdf,
    "md": parse_markdown,
}


def parse_document(file_bytes: bytes, file_type: str) -> List[ParsedChunk]:
    """Route to the correct parser based on file extension."""
    parser = PARSER_MAP.get(file_type)
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(file_bytes)